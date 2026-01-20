import os
import io
import json
import uuid
import zipfile
import tempfile
import subprocess
import fitz 
import base64

import io
import os
import tempfile
import requests
import difflib
from PyPDF2 import PdfReader
import io
from django.shortcuts import render
from django.http import HttpResponse
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.colors import black


from django.shortcuts import render
from django.http import HttpResponse

from PIL import Image
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4

import pdfkit

from django.shortcuts import render
from django.http import HttpResponse
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile

from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from django.conf import settings

import pdfkit
import tempfile
from django.http import HttpResponse
from django.shortcuts import render

WKHTML_PATH = r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe"
config = pdfkit.configuration(wkhtmltopdf=WKHTML_PATH)




GS_PATH = settings.GS_PATH


# ===============================
# HOME
# ===============================
def home(request):
    tools = [
        {"name": "Merge PDF", "desc": "Combine PDFs", "icon": "https://cdn.lordicon.com/ivhjpjsw.json", "url": "/merge/"},
        {"name": "Split PDF", "desc": "Split PDF pages", "icon": "https://cdn.lordicon.com/fhtaantg.json", "url": "/split/"},
        {"name": "Compress PDF", "desc": "Reduce PDF size", "icon": "https://cdn.lordicon.com/aklfruoc.json", "url": "/compress/"},
        {"name": "Rotate PDF", "desc": "Rotate pages", "icon": "https://cdn.lordicon.com/puvaffet.json", "url": "/rotate/"},
        {"name": "Protect PDF", "desc": "Add password", "icon": "https://cdn.lordicon.com/lxotnbfa.json", "url": "/protect/"},
        {"name": "Unlock PDF", "desc": "Remove password", "icon": "https://cdn.lordicon.com/tdrtiskw.json", "url": "/unlock/"},
        {"name": "Watermark", "desc": "Add watermark", "icon": "https://cdn.lordicon.com/ygvjgdmk.json", "url": "/watermark/"},
        {"name": "OCR PDF", "desc": "Scanned PDF to text", "icon": "https://cdn.lordicon.com/kkvxgpti.json", "url": "/ocr/"},
        {"name": "PDF → Word", "desc": "Convert PDF to DOCX", "icon": "https://cdn.lordicon.com/wzwygmng.json", "url": "/pdf-to-word/"},
        {"name": "Word → PDF", "desc": "Convert DOCX to PDF", "icon": "https://cdn.lordicon.com/qhgmphtg.json", "url": "/word-to-pdf/"},
        {"name": "PDF → Image", "desc": "PDF pages to images", "icon": "https://cdn.lordicon.com/rwotyanb.json", "url": "/pdf-to-image/"},
        {"name": "Image → PDF", "desc": "Images to PDF", "icon": "https://cdn.lordicon.com/iltqorsz.json", "url": "/image-to-pdf/"},
        {"name": "PDF → Excel", "desc": "Extract tables", "icon": "https://cdn.lordicon.com/puvaffet.json", "url": "/pdf-to-excel/"},
        {"name": "Excel → PDF", "desc": "Spreadsheet to PDF", "icon": "https://cdn.lordicon.com/jjoolpwc.json", "url": "/excel-to-pdf/"},
        {"name": "PPT → PDF", "desc": "Slides to PDF", "icon": "https://cdn.lordicon.com/iawrhwdo.json", "url": "/ppt-to-pdf/"},
        {"name": "PDF → PPT", "desc": "PDF to slides", "icon": "https://cdn.lordicon.com/nocovwne.json", "url": "/pdf-to-ppt/"},
        {
            "name": "Sign PDF","desc": "Secure digital signature","icon": "https://cdn.lordicon.com/xtzvywzp.json",
            "desc": "Secure digital signature",
            "icon": "https://cdn.lordicon.com/xtzvywzp.json",
            "url": "/sign/"
        },
        {
            "name": "Organize PDF",
            "desc": "Edit PDF structure",
            "icon": "https://cdn.lordicon.com/gsqxdxog.json",
            "url": "/organize/"
        },

        {
            "name": "Repair PDF",
            "desc": "Fix corrupted PDF files",
            "icon": "https://cdn.lordicon.com/wzwygmng.json",
            "url": "/repair/"
        },
        {
            "name": "Compare PDF",
            "desc": "Check PDF changes",
            "icon": "https://cdn.lordicon.com/ivhjpjsw.json",
            "url": "/compare/"
        },
        {
            "name": "Redact PDF",
            "desc": "Hide sensitive information",
            "icon": "https://cdn.lordicon.com/uqpazftn.json",
            "url": "/redact/"
        },
        {
            "name": "Scan to PDF",
            "desc": "Turn scans into PDFs",
            "icon": "https://cdn.lordicon.com/iltqorsz.json",
            "url": "/scan-to-pdf/"
        },
        {
            "name": "Page Numbers",
            "desc": "Add page numbers to PDF",
            "icon": "https://cdn.lordicon.com/puvaffet.json",
            "url": "/page-numbers/"
        },
        {
            "name": "HTML to PDF",
            "desc": "Convert web pages to PDF",
            "icon": "https://cdn.lordicon.com/xtzvywzp.json",
            "url": "/html-to-pdf/"
        },

        {
            "name": "PDF → HTML",
            "desc": "Convert PDF into web-friendly HTML",
            "icon": "https://cdn.lordicon.com/ylvuooxd.json",
            "url": "/pdf-to-html/"
        },
        {
            "name": "Crop PDF",
            "desc": "Trim PDF pages",
            "icon": "https://cdn.lordicon.com/aklfruoc.json",
            "url": "/crop/"
        },



    ]
    return render(request, "home.html", {"tools": tools})


# ===============================
# MERGE PDF
# ===============================
def merge_pdf(request):
    if request.method == "POST":
        merger = PdfMerger()
        for f in request.FILES.getlist("pdfs"):
            merger.append(f)

        buffer = io.BytesIO()
        merger.write(buffer)
        buffer.seek(0)

        return HttpResponse(buffer.read(), content_type="application/pdf")
    return render(request, "merge.html")


# ===============================
# SPLIT PDF
# ===============================
def split_pdf(request):
    if request.method == "POST":
        pdf = request.FILES.get("pdf")
        ranges = json.loads(request.POST.get("ranges", "[]"))

        reader = PdfReader(pdf)
        writer = PdfWriter()

        for r in ranges:
            for i in range(r["from"] - 1, r["to"]):
                writer.add_page(reader.pages[i])

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)
        return HttpResponse(buffer.read(), content_type="application/pdf")

    return render(request, "split.html")

import io
import os
from PyPDF2 import PdfReader, PdfWriter
from django.http import HttpResponse
from django.shortcuts import render

# ===============================
# COMPRESS PDF (Render-safe)
# ===============================
def compress_pdf(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf")
        level = request.POST.get("level", "recommended")

        if not pdf_file:
            return HttpResponse("No PDF uploaded", status=400)

        reader = PdfReader(pdf_file)
        writer = PdfWriter()

        # Compression levels
        compress = {
            "extreme": True,
            "recommended": True,
            "less": False
        }.get(level, True)

        for page in reader.pages:
            if compress:
                page.compress_content_streams()
            writer.add_page(page)

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)

        filename = os.path.splitext(pdf_file.name)[0] + "_compressed.pdf"

        response = HttpResponse(
            buffer.read(),
            content_type="application/pdf"
        )
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "compress.html")


# ===============================
# ROTATE PDF
# ===============================
def rotate_pdf(request):
    if request.method == "POST":
        pdf = request.FILES.get("pdf")
        rotation = int(request.POST.get("rotation", 90))

        reader = PdfReader(pdf)
        writer = PdfWriter()

        for p in reader.pages:
            p.rotate(rotation)
            writer.add_page(p)

        buf = io.BytesIO()
        writer.write(buf)
        buf.seek(0)
        return HttpResponse(buf.read(), content_type="application/pdf")

    return render(request, "rotate.html")


# ===============================
# WATERMARK
# ===============================
def watermark_pdf(request):
    if request.method == "POST":
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.colors import gray

        pdf = request.FILES.get("pdf")
        text = request.POST.get("text", "VetriFiles")

        reader = PdfReader(pdf)
        writer = PdfWriter()

        for page in reader.pages:
            packet = io.BytesIO()
            c = canvas.Canvas(packet, pagesize=A4)
            c.setFillColor(gray, alpha=0.2)
            c.setFont("Helvetica-Bold", 40)
            c.drawCentredString(300, 400, text)
            c.save()

            packet.seek(0)
            watermark = PdfReader(packet).pages[0]
            page.merge_page(watermark)
            writer.add_page(page)

        buf = io.BytesIO()
        writer.write(buf)
        buf.seek(0)
        return HttpResponse(buf.read(), content_type="application/pdf")

    return render(request, "watermark.html")


# ===============================
# OCR PDF
# ===============================
def ocr_pdf(request):
    if request.method == "POST":
        from pdf2image import convert_from_bytes
        import pytesseract

        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        POPPLER_PATH = r"C:\poppler\Library\bin"

        pdf = request.FILES.get("pdf")
        images = convert_from_bytes(pdf.read(), poppler_path=POPPLER_PATH)

        writer = PdfWriter()
        for img in images:
            b = pytesseract.image_to_pdf_or_hocr(img, extension="pdf")
            r = PdfReader(io.BytesIO(b))
            writer.add_page(r.pages[0])

        buf = io.BytesIO()
        writer.write(buf)
        buf.seek(0)
        return HttpResponse(buf.read(), content_type="application/pdf")

    return render(request, "ocr.html")


# ===============================
# PDF → WORD
# ===============================
def pdf_to_word(request):
    if request.method == "POST":
        from pdf2docx import Converter

        pdf = request.FILES.get("file")
        with tempfile.TemporaryDirectory() as tmp:
            pdf_p = os.path.join(tmp, "in.pdf")
            docx_p = os.path.join(tmp, "out.docx")
            open(pdf_p, "wb").write(pdf.read())

            cv = Converter(pdf_p)
            cv.convert(docx_p)
            cv.close()

            return HttpResponse(
                open(docx_p, "rb").read(),
                content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
    return render(request, "pdf_to_word.html")


# ===============================
# WORD → PDF
# ===============================
def word_to_pdf(request):
    if request.method == "POST":
        from docx2pdf import convert

        doc = request.FILES.get("file")
        with tempfile.TemporaryDirectory() as tmp:
            d = os.path.join(tmp, doc.name)
            p = d.replace(".docx", ".pdf")
            open(d, "wb").write(doc.read())
            convert(d, p)
            return HttpResponse(open(p, "rb").read(), content_type="application/pdf")

    return render(request, "word_to_pdf.html")


# ===============================
# PDF → IMAGE
# ===============================
def pdf_to_image(request):
    if request.method == "POST":
        import fitz

        pdf = request.FILES.get("pdf")
        doc = fitz.open(stream=pdf.read(), filetype="pdf")

        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, "w") as z:
            for i in range(len(doc)):
                pix = doc[i].get_pixmap()
                z.writestr(f"page_{i+1}.png", pix.tobytes("png"))

        zip_buf.seek(0)
        return HttpResponse(zip_buf.read(), content_type="application/zip")

    return render(request, "pdf_to_image.html")


# ===============================
# IMAGE → PDF
# ===============================
def image_to_pdf(request):
    from PIL import Image

    if request.method == "POST":
        imgs = [Image.open(f).convert("RGB") for f in request.FILES.getlist("images")]
        buf = io.BytesIO()
        imgs[0].save(buf, save_all=True, append_images=imgs[1:], format="PDF")
        buf.seek(0)
        return HttpResponse(buf.read(), content_type="application/pdf")

    return render(request, "image_to_pdf.html")


# ===============================
# PDF → EXCEL
# ===============================
def pdf_to_excel(request):
    if request.method == "POST":
        import pandas as pd
        import tabula

        pdf = request.FILES.get("file")
        with tempfile.TemporaryDirectory() as tmp:
            p = os.path.join(tmp, "in.pdf")
            x = os.path.join(tmp, "out.xlsx")
            open(p, "wb").write(pdf.read())

            tables = tabula.read_pdf(p, pages="all", multiple_tables=True)
            with pd.ExcelWriter(x) as w:
                for i, t in enumerate(tables):
                    t.to_excel(w, sheet_name=f"Sheet{i+1}", index=False)

            return HttpResponse(open(x, "rb").read(),
                content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    return render(request, "pdf_to_excel.html")


# ===============================
# EXCEL → PDF
# ===============================
def excel_to_pdf(request):
    if request.method == "POST":
        import pandas as pd
        import pdfkit

        excel = request.FILES.get("file")
        with tempfile.TemporaryDirectory() as tmp:
            e = os.path.join(tmp, excel.name)
            p = os.path.join(tmp, "out.pdf")
            open(e, "wb").write(excel.read())

            xl = pd.ExcelFile(e)
            html = "".join(xl.parse(s).to_html(index=False) for s in xl.sheet_names)
            pdfkit.from_string(html, p)

            return HttpResponse(open(p, "rb").read(), content_type="application/pdf")

    return render(request, "excel_to_pdf.html")


# ===============================
# PPT → PDF
# ===============================
def ppt_to_pdf(request):
    if request.method == "POST":
        ppt = request.FILES.get("file")
        with tempfile.TemporaryDirectory() as tmp:
            p = os.path.join(tmp, ppt.name)
            open(p, "wb").write(ppt.read())

            subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", p, "--outdir", tmp], check=True)
            return HttpResponse(open(p.replace(".pptx", ".pdf"), "rb").read(), content_type="application/pdf")

    return render(request, "ppt_to_pdf.html")


# ===============================
# PDF → PPT
# ===============================
def pdf_to_ppt(request):
    if request.method == "POST":
        from pdf2image import convert_from_bytes
        from pptx import Presentation

        pdf = request.FILES.get("file")
        prs = Presentation()
        imgs = convert_from_bytes(pdf.read(), dpi=150)

        for img in imgs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(buf, 0, 0, width=prs.slide_width, height=prs.slide_height)

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)

        return HttpResponse(
            out.read(),
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    return render(request, "pdf_to_ppt.html")

# ===============================   
# PROTECT PDF
# ===============================
def protect_pdf(request):
    if request.method == "POST":
        pdf = request.FILES.get("pdf")
        password = request.POST.get("password")

        reader = PdfReader(pdf)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(password)

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)

        return HttpResponse(
            buffer.read(),
            content_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=protected.pdf"
            }
        )

    return render(request, "protect.html")
# ===============================
# UNLOCK PDF
# ===============================
def unlock_pdf(request):
    if request.method == "POST":
        pdf = request.FILES.get("pdf")
        password = request.POST.get("password")

        reader = PdfReader(pdf)

        if reader.is_encrypted:
            try:
                reader.decrypt(password)
            except:
                return HttpResponse("Invalid password", status=400)

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)

        return HttpResponse(
            buffer.read(),
            content_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=unlocked.pdf"
            }
        )

    return render(request, "unlock.html")
import io
import os
import base64
import tempfile
from django.shortcuts import render
from django.http import HttpResponse
from PIL import Image
import fitz  # PyMuPDF
import pdfkit
import pdfplumber
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.colors import black

# ===============================
# SIGN PDF
# ===============================
def sign_pdf(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf")
        signature_data = request.POST.get("signature")

        if not pdf_file or not signature_data:
            return HttpResponse("Missing PDF or signature", status=400)

        # Decode signature
        header, encoded = signature_data.split(",", 1)
        signature_bytes = base64.b64decode(encoded)
        signature_img = Image.open(io.BytesIO(signature_bytes))

        # Open PDF
        pdf_bytes = pdf_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        page = doc[-1]

        # Convert PIL image to Pixmap
        img_bytes = io.BytesIO()
        signature_img.save(img_bytes, format="PNG")
        img_bytes.seek(0)
        img_pix = fitz.Pixmap(fitz.open("png", img_bytes))

        # Insert at bottom-right
        rect = fitz.Rect(page.rect.width - 200, page.rect.height - 120,
                         page.rect.width - 20, page.rect.height - 20)
        page.insert_image(rect, pixmap=img_pix)

        # Save to memory
        output = io.BytesIO()
        doc.save(output)
        doc.close()
        output.seek(0)

        # Download with original filename
        base_name = pdf_file.name.rsplit(".", 1)[0]
        filename = f"{base_name}_signed.pdf"

        response = HttpResponse(output.read(), content_type="application/pdf")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "sign_pdf.html")


# ===============================
# SCAN TO PDF
# ===============================
def scan_to_pdf(request):
    if request.method == "POST":
        images = request.FILES.getlist("images")
        if not images:
            return HttpResponse("No images uploaded", status=400)

        pil_images = [Image.open(img).convert("RGB") for img in images]

        buffer = io.BytesIO()
        pil_images[0].save(buffer, format="PDF", save_all=True, append_images=pil_images[1:])
        buffer.seek(0)

        # Filename based on first image if available
        base_name = images[0].name.rsplit(".", 1)[0] if images else "scanned"
        filename = f"{base_name}_scanned.pdf"

        response = HttpResponse(buffer.read(), content_type="application/pdf")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "scan_to_pdf.html")




# ===============================
# HTML TO PDF
# ===============================
from django.shortcuts import render
from django.http import HttpResponse
# from weasyprint import HTML


def html_to_pdf(request):
    if request.method == "POST":
        html_content = request.POST.get("html")
        filename = request.POST.get("filename", "document").strip() or "document"

        if not html_content:
            return HttpResponse("No HTML provided", status=400)

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            pdfkit.from_string(
                html_content,
                tmp.name,
                configuration=config,
                options={
                    "encoding": "UTF-8",
                    "enable-local-file-access": ""
                }
            )

        with open(tmp.name, "rb") as f:
            response = HttpResponse(f.read(), content_type="application/pdf")
            response["Content-Disposition"] = f'attachment; filename="{filename}.pdf"'
            return response

    return render(request, "html_to_pdf.html")






import io
import os
from django.shortcuts import render
from django.http import HttpResponse
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from django.views.decorators.csrf import csrf_exempt

# ===============================
# PAGE NUMBERS
# ===============================
@csrf_exempt
def page_numbers(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf")
        position = request.POST.get("position", "bottom")

        if not pdf_file:
            return HttpResponse("No PDF uploaded", status=400)

        reader = PdfReader(pdf_file)
        writer = PdfWriter()

        for i, page in enumerate(reader.pages, start=1):
            packet = io.BytesIO()
            c = canvas.Canvas(packet, pagesize=A4)
            c.setFont("Helvetica", 10)
            if position == "top":
                c.drawCentredString(300, 820, str(i))
            else:
                c.drawCentredString(300, 20, str(i))
            c.save()
            packet.seek(0)

            overlay = PdfReader(packet).pages[0]
            page.merge_page(overlay)
            writer.add_page(page)

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)

        base_name = os.path.splitext(pdf_file.name)[0]
        filename = f"{base_name}_pagenumbers.pdf"

        response = HttpResponse(buffer.read(), content_type="application/pdf")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "page_numbers.html")


# ===============================
# ORGANIZE PDF
# ===============================
@csrf_exempt
def organize_pdf(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf")
        order = request.POST.get("order")  # e.g. 3,1,2

        if not pdf_file or not order:
            return HttpResponse("Missing data", status=400)

        page_order = [int(x.strip()) - 1 for x in order.split(",")]

        reader = PdfReader(pdf_file)
        writer = PdfWriter()

        for i in page_order:
            if 0 <= i < len(reader.pages):
                writer.add_page(reader.pages[i])

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)

        base_name = os.path.splitext(pdf_file.name)[0]
        filename = f"{base_name}_organized.pdf"

        response = HttpResponse(buffer.read(), content_type="application/pdf")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "organize_pdf.html")


# ===============================
# COMPARE PDF
# ===============================
@csrf_exempt
def compare_pdf(request):
    if request.method == "POST":
        pdf1_file = request.FILES.get("pdf1")
        pdf2_file = request.FILES.get("pdf2")

        if not pdf1_file or not pdf2_file:
            return HttpResponse("Both PDF files are required.", status=400)

        pdf1 = PdfReader(pdf1_file)
        pdf2 = PdfReader(pdf2_file)
        writer = PdfWriter()

        max_pages = max(len(pdf1.pages), len(pdf2.pages))
        for i in range(max_pages):
            page1 = pdf1.pages[i] if i < len(pdf1.pages) else None
            page2 = pdf2.pages[i] if i < len(pdf2.pages) else None

            if page1 and page2 and page1.extract_text() == page2.extract_text():
                writer.add_page(page1)
            else:
                # Difference page
                packet = io.BytesIO()
                c = canvas.Canvas(packet, pagesize=A4)
                c.drawString(100, 500, f"Page {i+1}: Difference detected")
                c.save()
                packet.seek(0)
                diff_page = PdfReader(packet).pages[0]
                writer.add_page(diff_page)

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)

        base_name = os.path.splitext(pdf1_file.name)[0]
        filename = f"{base_name}_comparison.pdf"

        response = HttpResponse(buffer.read(), content_type="application/pdf")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "compare_pdf.html")



# ===============================
# REPAIR PDF
# ===============================
def repair_pdf(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf")
        if not pdf_file:
            return HttpResponse("No PDF uploaded", status=400)

        reader = PdfReader(pdf_file, strict=False)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)

        base_name = pdf_file.name.rsplit(".", 1)[0]
        filename = f"{base_name}_repaired.pdf"

        response = HttpResponse(buffer.read(), content_type="application/pdf")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "repair_pdf.html")





# ===============================
# REDACT PDF
# ===============================
def redact_pdf(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf")
        text_to_redact = request.POST.get("text", "")

        if not pdf_file:
            return HttpResponse("No PDF uploaded", status=400)

        reader = PdfReader(pdf_file)
        writer = PdfWriter()

        for page in reader.pages:
            packet = io.BytesIO()
            c = canvas.Canvas(packet, pagesize=A4)
            c.setFillColor(black)
            c.rect(50, 50, 500, 50, fill=1)  # You can adjust redaction box
            c.save()
            packet.seek(0)

            overlay = PdfReader(packet).pages[0]
            page.merge_page(overlay)
            writer.add_page(page)

        buffer = io.BytesIO()
        writer.write(buffer)
        buffer.seek(0)

        base_name = pdf_file.name.rsplit(".", 1)[0]
        filename = f"{base_name}_redacted.pdf"

        response = HttpResponse(buffer.read(), content_type="application/pdf")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "redact_pdf.html")


# ===============================
# PDF TO HTML
# ===============================
def pdf_to_html(request):
    if request.method == "POST":
        pdf = request.FILES.get("pdf")
        if not pdf:
            return HttpResponse("No PDF uploaded", status=400)

        base_name = os.path.splitext(pdf.name)[0]
        output_filename = f"{base_name}.html"

        with tempfile.TemporaryDirectory() as tmp:
            pdf_path = os.path.join(tmp, pdf.name)
            with open(pdf_path, "wb") as f:
                for chunk in pdf.chunks():
                    f.write(chunk)

            html = ["<!DOCTYPE html>", "<html>", "<head>", "<meta charset='utf-8'>",
                    f"<title>{base_name}</title>", "<style>body{font-family:Arial; line-height:1.6; padding:20px} p{margin-bottom:12px}</style>",
                    "</head><body>"]

            import pdfplumber
            with pdfplumber.open(pdf_path) as pdf_file:
                for i, page in enumerate(pdf_file.pages, start=1):
                    text = page.extract_text()
                    if text:
                        html.append(f"<h3>Page {i}</h3>")
                        html.append(f"<p>{text.replace(chr(10), '<br>')}</p>")
            html.append("</body></html>")

        response = HttpResponse("\n".join(html), content_type="text/html")
        response["Content-Disposition"] = f'attachment; filename="{output_filename}"'
        return response

    return render(request, "pdf_to_html.html")
#==============================
# Crop PDF 
#=============================
import io
from django.shortcuts import render
from django.http import HttpResponse
import fitz  # PyMuPDF

def crop_pdf(request):
    """
    Upload PDF, crop via mouse or manual input, and download with _crop filename.
    """
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf")
        if not pdf_file:
            return HttpResponse("No PDF uploaded", status=400)

        # Safely parse coordinates
        try:
            left = float(request.POST.get("left") or 0)
            top = float(request.POST.get("top") or 0)
            right = float(request.POST.get("right") or 0)
            bottom = float(request.POST.get("bottom") or 0)
        except ValueError:
            return HttpResponse("Invalid crop values", status=400)

        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        for page in doc:
            rect = page.rect
            crop_rect = fitz.Rect(left, top, right, bottom) & rect
            page.set_cropbox(crop_rect)

        buffer = doc.write()
        response = HttpResponse(buffer, content_type="application/pdf")
        filename = pdf_file.name.replace(".pdf","_crop.pdf")
        response["Content-Disposition"] = f'attachment; filename="{filename}"'
        return response

    return render(request, "crop_pdf.html")
